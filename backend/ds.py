import os

from bs4 import BeautifulSoup
from docx import Document
from odf.opendocument import load
from odf.text import P
from pptx import Presentation 
from pypdf import PdfReader
from striprtf.striprtf import rtf_to_text
from tika import parser

from transformers import BartTokenizer, BartForConditionalGeneration


tokenizer = BartTokenizer.from_pretrained('facebook/bart-large-cnn')
model = BartForConditionalGeneration.from_pretrained('facebook/bart-large-cnn') 

def extract_text(filepath) -> str:

    ext = os.path.splitext(filepath)[1].lstrip('.')
    text : str = ""

    match ext:
        
        case "docx":
            doc = Document(filepath)
            text = "\n".join(p.text for p in doc.paragraphs)

        case "html":
            with open("file.html", "r", encoding="utf-8", errors="ignore") as f:
                soup = BeautifulSoup(f, "lxml")

            text = soup.get_text(separator="\n", strip=True)

        case "odt":
            doc = load("file.odt")
            paras = doc.getElementsByType(P)
            text = "\n".join(p.firstChild.data for p in paras if p.firstChild)

        case "pdf": 
            reader = PdfReader(filepath)
    
            for page in reader.pages:
                pdftext = page.extract_text()
                if pdftext:                 # Check if text successfully extracted
                    text += pdftext + "\n"  # Add newline between pages

        case "pptx":
            prs = Presentation("file.pptx")
            text_runs = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape)

            text = "\n".join(text_runs)

        case "rtf":
            with open("file.rtf", "r", encoding="utf-8", errors="ignore") as f:
                text = rtf_to_text(f.read())

        case "txt":
            with open("file.txt", "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        case _ :
            return text

    return text


def generate_summary(filepath, instruction):

    text : str = extract_text(filepath)

    inputs = tokenizer.encode(instruction + "\nSummarize: " + text, return_tensors="pt", max_length=1024, truncation=True)
    summary_ids = model.generate(inputs, max_length=150, min_length=50, length_penalty=2.0, num_beams=4, early_stopping=True)
    summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
    return summary